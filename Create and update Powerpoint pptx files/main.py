'''
Creating Presentations with Python
https://towardsdatascience.com/creating-presentations-with-python-3f5737824f61?gi=c6666401f8cb

One survey, 100 decks and 1,000 slides in 10 minutes
https://towardsdatascience.com/one-survey-100-decks-and-1-000-slides-in-10-minutes-8fcfd1b246ee

pip install python-pptx

pip install plotly

npm install -g electron@6.1.4 orca
'''

from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import ChartData


prs=Presentation()
lyt=prs.slide_layouts[0] # choosing a slide layout
slide=prs.slides.add_slide(lyt) # adding a slide
title=slide.shapes.title # assigning a title
subtitle=slide.placeholders[1] # placeholder for subtitle
title.text="Hey,This is a Slide! How exciting!" # title
subtitle.text="Really?" # subtitle
# prs.save("slide1.pptx") # saving file


# If you want to use a widescreen presentation, you should define slide sizes in the beginning
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# prs = Presentation()
blank_slide_layout = prs.slide_layouts[6] # 6 is a blank slide
slide = prs.slides.add_slide(blank_slide_layout)

# # Adding Images
path='stick2.png' 
left=Inches(1)
top=Inches(0.5)
img=slide.shapes.add_picture(path,left,top)

# Adding Chart
## Let’s add some categories and series.
prs=Presentation()
# define some measures. (left, top, width, height)
x,y,cx,cy=Inches(1),Inches(2),Inches(5),Inches(7)
## add chart to slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

chart_data=ChartData()
chart_data.categories = ['Column 1', 'Column 2', 'Column 3']
chart_data.add_series('Hey',    (10.5, 5.5, 17.5))
chart_data.add_series('there',    (25.5, 40.3, 30.7))
chart_data.add_series('reader', (5.2, 10.3, 8.4))
chart = slide.shapes.add_chart( XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data).chart

# Adding more blank slide
for _ in range(10):
    prs.slides.add_slide(prs.slide_layouts[6])

prs.save("blankslide.pptx") # saving file


